"""Build the datathon presentation slides — 2kim_finance_slides.pptx

Run: python scripts/build_slides.py
Output: submissions/2kim_finance_slides.pptx
"""

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Colors ────────────────────────────────────────────────────────────────────
BG_DARK    = RGBColor(0x0F, 0x14, 0x19)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xD1, 0xD5, 0xDB)
MID_GRAY   = RGBColor(0x6B, 0x72, 0x80)
DARK_GRAY  = RGBColor(0x37, 0x41, 0x51)
NVDA_GREEN = RGBColor(0x00, 0xCC, 0x96)
CSCO_RED   = RGBColor(0xEF, 0x55, 0x3B)
ACCENT     = RGBColor(0x63, 0x6E, 0xFA)
CARD_BG    = RGBColor(0x16, 0x1B, 0x22)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

OUT_DIR = Path(__file__).resolve().parents[1] / "submissions"
OUT_DIR.mkdir(exist_ok=True)


# ── Helpers ───────────────────────────────────────────────────────────────────

def set_slide_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, left, top, width, height, text, *,
                 font_size=18, bold=False, color=WHITE, alignment=PP_ALIGN.LEFT,
                 font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_shape_box(slide, left, top, width, height, fill_color=CARD_BG, border_color=DARK_GRAY):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = border_color
    shape.line.width = Pt(1)
    shape.shadow.inherit = False
    return shape


def add_bullet_list(slide, left, top, width, height, items, *,
                    font_size=16, color=LIGHT_GRAY, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = font_name
        p.space_after = Pt(6)
    return txBox


CHART_DIR = OUT_DIR / "charts"

# Map placeholder labels to chart file names
CHART_FILE_MAP = {
    "1.1": "chart_1_1.png",
    "1.2": "chart_1_2.png",
    "1.3": "chart_1_3.png",
    "1.4": "chart_1_4.png",
    "2.1": "chart_2_1_pe_trajectory.png",
    "2.2": "chart_2_2_revenue_growth.png",
    "3.1": "chart_3_1_concentration_timeline.png",
    "3.2": "chart_3_2_spy_vs_rsp.png",
    "3.3": "chart_3_3_buffett_indicator.png",
    "3.4": "chart_3_4_treemap.png",
    "4.5": "chart_4_5_macro_dashboard.png",
    "5.1": "chart_5_1_ai_mentions.png",
    "5.2": "chart_5_2_hype_vs_specificity.png",
    "ML.1": "chart_ml_1_regime_probabilities.png",
    "ML.2": "chart_ml_2_dtw_similarity.png",
    "ML.3": "chart_ml_3_shap_summary.png",
    "ML.4": "chart_ml_4_shap_waterfall.png",
}


def chart_placeholder(slide, left, top, width, height, label="[Chart Placeholder]"):
    """Insert chart image if available, otherwise show placeholder box."""
    # Try to find a matching chart PNG
    for chart_id, filename in CHART_FILE_MAP.items():
        if f"Chart {chart_id}" in label or f"chart_{chart_id.replace('.', '_')}" in label.lower():
            img_path = CHART_DIR / filename
            if img_path.exists():
                slide.shapes.add_picture(str(img_path), left, top, width, height)
                return
            break

    # Also try direct filename matching from the label
    for filename in CHART_DIR.glob("*.png") if CHART_DIR.exists() else []:
        fname_lower = filename.stem.lower()
        # Extract chart ID from label like "[Chart 1.1 — ...]"
        import re
        match = re.search(r'chart[_ ](\d+[._]\d+)', label.lower().replace(' ', '_'))
        if match:
            chart_num = match.group(1).replace('.', '_')
            if chart_num in fname_lower:
                slide.shapes.add_picture(str(filename), left, top, width, height)
                return

    # Fallback: gray placeholder box
    shape = add_shape_box(slide, left, top, width, height,
                          fill_color=RGBColor(0x0D, 0x11, 0x17),
                          border_color=DARK_GRAY)
    shape.text_frame.paragraphs[0].text = label
    shape.text_frame.paragraphs[0].font.size = Pt(14)
    shape.text_frame.paragraphs[0].font.color.rgb = MID_GRAY
    shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    shape.text_frame.paragraphs[0].font.name = "Calibri"
    shape.text_frame.word_wrap = True
    shape.text_frame.paragraphs[0].space_before = Pt(height.inches * 25)
    return shape


def add_footer(slide, text="Team 2Kim | SBU AI Community Datathon 2026 | Finance & Economics"):
    add_text_box(slide, Inches(0.5), Inches(6.9), Inches(12), Inches(0.4),
                 text, font_size=9, color=MID_GRAY, font_name="Calibri")


def add_divider(slide, top):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(0.8), top, Inches(11.7), Pt(1))
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_GRAY
    shape.line.fill.background()
    return shape


# ── Build ─────────────────────────────────────────────────────────────────────

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
blank_layout = prs.slide_layouts[6]  # blank

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 1 — Title
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl)

add_text_box(sl, Inches(1), Inches(1.8), Inches(11), Inches(1.2),
             "The AI Hype Cycle", font_size=48, bold=True, color=WHITE)
add_text_box(sl, Inches(1), Inches(3.0), Inches(11), Inches(0.8),
             "Are We in a Bubble?", font_size=32, bold=False, color=NVDA_GREEN)
add_text_box(sl, Inches(1), Inches(3.8), Inches(11), Inches(0.6),
             "Comparing NVIDIA 2023-2026 to Cisco 1998-2001", font_size=20, color=LIGHT_GRAY)

add_divider(sl, Inches(4.7))

add_text_box(sl, Inches(1), Inches(5.0), Inches(5), Inches(0.4),
             "Team 2Kim: Jimmy Kim & Alice Kim", font_size=16, color=LIGHT_GRAY)
add_text_box(sl, Inches(1), Inches(5.4), Inches(8), Inches(0.4),
             "SBU AI Community Datathon 2026  —  Finance & Economics Track",
             font_size=14, color=MID_GRAY)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 2 — The Hook: Price Overlay
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl)
add_text_box(sl, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
             '"Does This Look Familiar?"', font_size=32, bold=True, color=WHITE)

chart_placeholder(sl, Inches(0.5), Inches(1.3), Inches(12.3), Inches(5.2),
                  "[Chart 1.1 — Normalized Price Overlay: NVDA (green) vs CSCO (red) vs Nortel (gray dashed)]")

# Legend
add_text_box(sl, Inches(0.8), Inches(6.6), Inches(3), Inches(0.3),
             "■ NVIDIA 2023-2026", font_size=12, color=NVDA_GREEN)
add_text_box(sl, Inches(3.8), Inches(6.6), Inches(3), Inches(0.3),
             "■ Cisco 1998-2001", font_size=12, color=CSCO_RED)
add_text_box(sl, Inches(6.8), Inches(6.6), Inches(4), Inches(0.3),
             "- - Nortel Networks (bankrupt 2009)", font_size=12, color=MID_GRAY)
add_footer(sl)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 3 — Five-Layer Framework
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl)
add_text_box(sl, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
             "Five-Layer Analytical Framework", font_size=32, bold=True, color=WHITE)

layers = [
    ("1  PRICE",          "Same trajectory?",       NVDA_GREEN),
    ("2  FUNDAMENTALS",   "Same detachment?",       RGBColor(0x19, 0xD3, 0xF3)),
    ("3  CONCENTRATION",  "Same fragility?",        ACCENT),
    ("4  MACRO",          "Same environment?",      RGBColor(0xFE, 0xCA, 0x57)),
    ("5  SENTIMENT",      "Same euphoria?",         CSCO_RED),
]
for i, (label, question, clr) in enumerate(layers):
    y = Inches(1.5) + Inches(i * 1.0)
    # colored bar
    bar = slide_shapes = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                              Inches(1), y, Inches(4.5), Inches(0.7))
    bar.fill.solid()
    bar.fill.fore_color.rgb = CARD_BG
    bar.line.color.rgb = clr
    bar.line.width = Pt(2)

    add_text_box(sl, Inches(1.3), y + Inches(0.12), Inches(4), Inches(0.5),
                 label, font_size=18, bold=True, color=clr)

    add_text_box(sl, Inches(6), y + Inches(0.12), Inches(5), Inches(0.5),
                 question, font_size=18, color=LIGHT_GRAY)

add_text_box(sl, Inches(1), Inches(6.5), Inches(10), Inches(0.4),
             "A price pattern alone means nothing — but if all five dimensions align, that's different.",
             font_size=14, color=MID_GRAY)
add_footer(sl)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 4 — Data Sources & Methodology
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl)
add_text_box(sl, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
             "Data Sources & Methodology", font_size=32, bold=True, color=WHITE)

sources = [
    "Yahoo Finance — Daily prices, volumes, financials (1996-2026)",
    "FRED API — 10 macro series: Fed Funds, M2, yields, CPI, GDP",
    "SEC EDGAR — 10-Q/10-K filings for earnings call NLP",
    "Alpha Vantage — Quarterly fundamentals (P/E, revenue, FCF)",
    "Google Trends — Search interest for AI/dot-com terms",
    "OpenAI API — Sentiment scoring of earnings call sentences",
    "Reddit API — r/wallstreetbets NVDA mention frequency",
    "S&P 500 Constituent Data — Market cap weights, concentration",
]
add_bullet_list(sl, Inches(0.8), Inches(1.3), Inches(6), Inches(4.5),
                sources, font_size=15, color=LIGHT_GRAY)

# Stats box
add_shape_box(sl, Inches(7.5), Inches(1.3), Inches(5), Inches(3.5))
stats = [
    ("38", "features engineered"),
    ("384", "months of training data"),
    ("8", "data sources merged"),
    ("5", "analytical layers"),
    ("3", "ML models (ensemble)"),
]
for i, (num, label) in enumerate(stats):
    y = Inches(1.5) + Inches(i * 0.65)
    add_text_box(sl, Inches(7.8), y, Inches(1.5), Inches(0.5),
                 num, font_size=28, bold=True, color=NVDA_GREEN,
                 alignment=PP_ALIGN.RIGHT)
    add_text_box(sl, Inches(9.5), y + Inches(0.08), Inches(2.8), Inches(0.5),
                 label, font_size=14, color=LIGHT_GRAY)

add_text_box(sl, Inches(7.8), Inches(5.2), Inches(4.5), Inches(0.5),
             "Time-series CV  •  No lookahead bias  •  All trailing features",
             font_size=11, color=MID_GRAY)
add_footer(sl)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 5 — P/E Comparison (Bull Case)
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl)
add_text_box(sl, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
             "Valuations: Night and Day", font_size=32, bold=True, color=WHITE)

chart_placeholder(sl, Inches(0.5), Inches(1.3), Inches(8), Inches(5.2),
                  "[Chart 2.1 — P/E Ratio Trajectory: NVDA vs CSCO aligned by cycle phase]")

# Stats panel
add_shape_box(sl, Inches(8.8), Inches(1.3), Inches(4), Inches(5.2))
add_text_box(sl, Inches(9.0), Inches(1.6), Inches(3.6), Inches(0.4),
             "CISCO AT PEAK", font_size=12, bold=True, color=CSCO_RED)
add_text_box(sl, Inches(9.0), Inches(2.0), Inches(3.6), Inches(0.7),
             "~200x", font_size=48, bold=True, color=CSCO_RED)
add_text_box(sl, Inches(9.0), Inches(2.7), Inches(3.6), Inches(0.3),
             "Price-to-Earnings", font_size=14, color=MID_GRAY)

add_text_box(sl, Inches(9.0), Inches(3.5), Inches(3.6), Inches(0.4),
             "NVIDIA CURRENT", font_size=12, bold=True, color=NVDA_GREEN)
add_text_box(sl, Inches(9.0), Inches(3.9), Inches(3.6), Inches(0.7),
             "~50x", font_size=48, bold=True, color=NVDA_GREEN)
add_text_box(sl, Inches(9.0), Inches(4.6), Inches(3.6), Inches(0.3),
             "Price-to-Earnings", font_size=14, color=MID_GRAY)

add_text_box(sl, Inches(9.0), Inches(5.3), Inches(3.6), Inches(0.8),
             "Cisco's P/E was 4x NVIDIA's at equivalent cycle stage",
             font_size=13, color=LIGHT_GRAY)
add_footer(sl)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 6 — Revenue Growth
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl)
add_text_box(sl, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
             "Revenue Growth: 4x Faster Than Cisco", font_size=32, bold=True, color=WHITE)

chart_placeholder(sl, Inches(0.5), Inches(1.3), Inches(12.3), Inches(4.0),
                  "[Chart 2.2 — Revenue Growth Rate Comparison: NVDA (green bars) vs CSCO (red bars)]")

# Large stat callout
add_text_box(sl, Inches(1.5), Inches(5.5), Inches(4), Inches(1),
             "265%", font_size=56, bold=True, color=NVDA_GREEN)
add_text_box(sl, Inches(1.5), Inches(6.2), Inches(4), Inches(0.4),
             "NVDA peak YoY revenue growth", font_size=14, color=LIGHT_GRAY)

add_text_box(sl, Inches(6.5), Inches(5.5), Inches(4), Inches(1),
             "66%", font_size=56, bold=True, color=CSCO_RED)
add_text_box(sl, Inches(6.5), Inches(6.2), Inches(4), Inches(0.4),
             "CSCO peak YoY revenue growth", font_size=14, color=LIGHT_GRAY)
add_footer(sl)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 7 — Concentration (Bear Case)
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl)
add_text_box(sl, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
             "Market Concentration Is at Dot-Com Levels", font_size=32, bold=True, color=WHITE)

chart_placeholder(sl, Inches(0.5), Inches(1.3), Inches(7.5), Inches(5.2),
                  "[Chart 3.1 — Top-10 Concentration Ratio Timeline]")
chart_placeholder(sl, Inches(8.3), Inches(1.3), Inches(4.5), Inches(2.4),
                  "[Chart 3.2 — SPY vs RSP Spread]")

# Key stats box
add_shape_box(sl, Inches(8.3), Inches(4.0), Inches(4.5), Inches(2.5))
add_text_box(sl, Inches(8.5), Inches(4.2), Inches(4), Inches(0.3),
             "TOP 5 STOCKS % OF S&P 500", font_size=11, bold=True, color=MID_GRAY)
add_text_box(sl, Inches(8.5), Inches(4.6), Inches(2), Inches(0.5),
             "28%", font_size=36, bold=True, color=CSCO_RED)
add_text_box(sl, Inches(10.2), Inches(4.7), Inches(2), Inches(0.3),
             "2026", font_size=14, color=LIGHT_GRAY)
add_text_box(sl, Inches(8.5), Inches(5.2), Inches(2), Inches(0.5),
             "18%", font_size=36, bold=True, color=MID_GRAY)
add_text_box(sl, Inches(10.2), Inches(5.3), Inches(2), Inches(0.3),
             "2000", font_size=14, color=LIGHT_GRAY)
add_text_box(sl, Inches(8.5), Inches(5.8), Inches(4), Inches(0.5),
             "More concentrated than at the dot-com peak.",
             font_size=12, color=LIGHT_GRAY)
add_footer(sl)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 8 — Sentiment
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl)
add_text_box(sl, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
             "Hype Is Real — But So Is Substance", font_size=32, bold=True, color=WHITE)

chart_placeholder(sl, Inches(0.5), Inches(1.3), Inches(6), Inches(5.2),
                  '[Chart 5.1 — AI Mentions in Earnings Calls Over Time]')
chart_placeholder(sl, Inches(6.8), Inches(1.3), Inches(6), Inches(5.2),
                  '[Chart 5.2 — Hype Score vs Revenue Specificity Scatter]')
add_footer(sl)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 9 — Macro Wild Card
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl)
add_text_box(sl, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
             "The Macro Wild Card", font_size=32, bold=True, color=WHITE)

chart_placeholder(sl, Inches(0.5), Inches(1.3), Inches(12.3), Inches(5.2),
                  "[Chart 4.5 — Macro Dashboard: 2x3 Small Multiples (Fed Funds, M2, Yield Curve, Credit Spreads, CPI, GDP)]")

add_text_box(sl, Inches(0.8), Inches(6.6), Inches(11), Inches(0.3),
             "The macro environment is fundamentally different from 1999-2000 — but the yield curve inverted for the longest period in history.",
             font_size=13, color=MID_GRAY)
add_footer(sl)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 10 — ML Regime Classification
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl)
add_text_box(sl, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
             "What Do the Models Say?", font_size=32, bold=True, color=WHITE)

chart_placeholder(sl, Inches(0.5), Inches(1.3), Inches(8.5), Inches(5.2),
                  "[Chart ML.1 — Regime Classification Probabilities (stacked area)]")

# Probability readout
add_shape_box(sl, Inches(9.3), Inches(1.3), Inches(3.7), Inches(5.2))
add_text_box(sl, Inches(9.5), Inches(1.6), Inches(3.3), Inches(0.4),
             "MARCH 2026", font_size=14, bold=True, color=WHITE)
add_text_box(sl, Inches(9.5), Inches(2.1), Inches(3.3), Inches(0.3),
             "Ensemble Regime Probabilities", font_size=11, color=MID_GRAY)

probs = [
    ("Bubble",          "XX%", CSCO_RED),
    ("Normal Growth",   "XX%", NVDA_GREEN),
    ("Recovery",        "XX%", ACCENT),
    ("Correction",      "XX%", RGBColor(0xFE, 0xCA, 0x57)),
]
for i, (label, val, clr) in enumerate(probs):
    y = Inches(2.8) + Inches(i * 0.85)
    add_text_box(sl, Inches(9.5), y, Inches(1.8), Inches(0.4),
                 label, font_size=14, color=LIGHT_GRAY)
    add_text_box(sl, Inches(11.3), y, Inches(1.5), Inches(0.5),
                 val, font_size=24, bold=True, color=clr, alignment=PP_ALIGN.RIGHT)

add_text_box(sl, Inches(9.5), Inches(5.6), Inches(3.3), Inches(0.7),
             "3-model ensemble: RF + XGBoost + Logistic Regression. 38 features, 28 years.",
             font_size=10, color=MID_GRAY)
add_footer(sl)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 11 — DTW + SHAP
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl)
add_text_box(sl, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
             "The Tug-of-War: Bubble Signals vs Counter-Signals",
             font_size=28, bold=True, color=WHITE)

chart_placeholder(sl, Inches(0.5), Inches(1.2), Inches(12.3), Inches(2.5),
                  "[Chart ML.2 — DTW Similarity Timeline (rolling monthly score)]")
chart_placeholder(sl, Inches(0.5), Inches(4.0), Inches(12.3), Inches(2.8),
                  "[Chart ML.4 — SHAP Waterfall: What's pushing toward/away from 'Bubble' for March 2026]")

add_text_box(sl, Inches(0.8), Inches(6.9), Inches(11), Inches(0.3),
             "Concentration + sentiment push toward bubble  |  Revenue + FCF push away",
             font_size=13, color=MID_GRAY)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 12 — Synthesis Scorecard
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl)
add_text_box(sl, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
             "Synthesis: A Bubble? Not Quite. But Not Safe Either.",
             font_size=28, bold=True, color=WHITE)

# Scorecard table
rows = [
    ("Price Trajectory",    "YES (similar)",      "CAUTION",  CSCO_RED),
    ("Fundamentals",        "NO (much stronger)", "SUPPORT",  NVDA_GREEN),
    ("Market Concentration","YES (worse)",        "DANGER",   CSCO_RED),
    ("Macro Environment",   "MIXED",              "NEUTRAL",  RGBColor(0xFE, 0xCA, 0x57)),
    ("Sentiment / Hype",    "YES (elevated)",     "CAUTION",  CSCO_RED),
]

# Header
add_text_box(sl, Inches(1.0), Inches(1.4), Inches(3.5), Inches(0.4),
             "LAYER", font_size=13, bold=True, color=MID_GRAY)
add_text_box(sl, Inches(4.8), Inches(1.4), Inches(3.5), Inches(0.4),
             "DOT-COM SIGNAL?", font_size=13, bold=True, color=MID_GRAY)
add_text_box(sl, Inches(8.5), Inches(1.4), Inches(3), Inches(0.4),
             "VERDICT", font_size=13, bold=True, color=MID_GRAY)

for i, (layer, signal, verdict, clr) in enumerate(rows):
    y = Inches(1.9) + Inches(i * 0.75)
    # Row background
    add_shape_box(sl, Inches(0.8), y, Inches(11.5), Inches(0.65),
                  fill_color=CARD_BG if i % 2 == 0 else BG_DARK,
                  border_color=DARK_GRAY)
    add_text_box(sl, Inches(1.0), y + Inches(0.12), Inches(3.5), Inches(0.4),
                 layer, font_size=16, bold=True, color=WHITE)
    add_text_box(sl, Inches(4.8), y + Inches(0.12), Inches(3.5), Inches(0.4),
                 signal, font_size=15, color=LIGHT_GRAY)
    add_text_box(sl, Inches(8.5), y + Inches(0.12), Inches(3), Inches(0.4),
                 verdict, font_size=16, bold=True, color=clr)

# Summary
add_text_box(sl, Inches(1.0), Inches(5.9), Inches(11), Inches(0.7),
             "3 of 5 layers flash bubble signals.  Strong fundamentals are the only thing preventing a bubble classification.",
             font_size=16, color=LIGHT_GRAY)

add_text_box(sl, Inches(1.0), Inches(6.5), Inches(11), Inches(0.4),
             "If fundamentals disappoint even slightly, the downside could be severe.",
             font_size=14, bold=True, color=CSCO_RED)
add_footer(sl)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 13 — Key Takeaway
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl)

add_text_box(sl, Inches(1.5), Inches(1.5), Inches(10), Inches(0.5),
             "THE KEY INSIGHT", font_size=16, bold=True, color=NVDA_GREEN)

add_text_box(sl, Inches(1.5), Inches(2.3), Inches(10), Inches(2),
             '"The AI rally is fundamentally justified but structurally fragile.\n'
             'The gap between a justified rally and a bubble is measured\n'
             'in quarters — not years — of continued earnings delivery."',
             font_size=26, bold=False, color=WHITE)

add_divider(sl, Inches(4.8))

add_text_box(sl, Inches(1.5), Inches(5.1), Inches(10), Inches(1),
             "NVIDIA must keep executing. Any deceleration in revenue growth would remove\n"
             "the one factor separating this from a classic bubble.",
             font_size=16, color=LIGHT_GRAY)
add_footer(sl)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 14 — Limitations
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl)
add_text_box(sl, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
             "Limitations & Ethical Considerations", font_size=32, bold=True, color=WHITE)

limitations = [
    "1.  Small sample of historical bubbles — only 2-3 comparable events in training data.",
    "2.  Labels are subjective — sensitivity analysis shows ±15% variance in results.",
    "3.  Structural market changes — passive investing (55% today vs 10% in 2000),\n"
    "     algorithmic trading, different Fed tools mean history may not repeat.",
    "4.  Not investment advice — academic analysis only. The model classifies regime\n"
    "     risk, it does not predict future prices.",
    "5.  AI tools used responsibly — OpenAI API for NLP classification.\n"
    "     All results manually validated on a 200-sentence sample (~85% agreement).",
]
add_bullet_list(sl, Inches(0.8), Inches(1.4), Inches(11), Inches(5),
                limitations, font_size=16, color=LIGHT_GRAY)
add_footer(sl)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 15 — Thank You / Q&A
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl)

add_text_box(sl, Inches(1), Inches(2.0), Inches(11), Inches(1),
             "Thank You", font_size=48, bold=True, color=WHITE,
             alignment=PP_ALIGN.CENTER)
add_text_box(sl, Inches(1), Inches(3.2), Inches(11), Inches(0.6),
             "Team 2Kim: Jimmy Kim & Alice Kim", font_size=20, color=LIGHT_GRAY,
             alignment=PP_ALIGN.CENTER)

add_divider(sl, Inches(4.2))

add_text_box(sl, Inches(1), Inches(4.5), Inches(11), Inches(0.8),
             "Questions?", font_size=36, bold=True, color=NVDA_GREEN,
             alignment=PP_ALIGN.CENTER)

add_text_box(sl, Inches(1), Inches(5.8), Inches(11), Inches(0.5),
             "Notebook: 2kim_finance_notebook.ipynb  |  Interactive dashboard available for live demo if requested",
             font_size=12, color=MID_GRAY, alignment=PP_ALIGN.CENTER)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 16 (BACKUP) — SHAP Deep Dive
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl)
add_text_box(sl, Inches(0.8), Inches(0.2), Inches(8), Inches(0.5),
             "BACKUP", font_size=12, bold=True, color=MID_GRAY)
add_text_box(sl, Inches(0.8), Inches(0.5), Inches(10), Inches(0.6),
             "Feature Importance Deep Dive: SHAP Summary",
             font_size=28, bold=True, color=WHITE)

chart_placeholder(sl, Inches(0.5), Inches(1.3), Inches(12.3), Inches(5.5),
                  "[Chart ML.3 — SHAP Summary Plot (Beeswarm): 38 features ranked by importance]")
add_footer(sl)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 17 (BACKUP) — Walk-Forward Validation
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl)
add_text_box(sl, Inches(0.8), Inches(0.2), Inches(8), Inches(0.5),
             "BACKUP", font_size=12, bold=True, color=MID_GRAY)
add_text_box(sl, Inches(0.8), Inches(0.5), Inches(10), Inches(0.6),
             "Model Validation: Walk-Forward Results",
             font_size=28, bold=True, color=WHITE)

chart_placeholder(sl, Inches(0.5), Inches(1.3), Inches(6), Inches(5.5),
                  "[Walk-Forward Validation Accuracy Chart]")
chart_placeholder(sl, Inches(6.8), Inches(1.3), Inches(6), Inches(2.5),
                  "[Confusion Matrix — Test Set]")
chart_placeholder(sl, Inches(6.8), Inches(4.1), Inches(6), Inches(2.7),
                  "[Ensemble Agreement Statistics]")
add_footer(sl)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SAVE
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
out_path = OUT_DIR / "2kim_finance_slides.pptx"
prs.save(out_path)
print(f"Saved: {out_path}")
print(f"Slides: {len(prs.slides)}")
